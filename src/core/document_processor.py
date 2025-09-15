import os
from pathlib import Path
from typing import List, Dict
from langchain.schema import Document

class DocumentProcessor:
    """文档处理器，支持多种格式包括Word/WPS"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.docx': self._process_word,
            '.doc': self._process_word,
            '.wps': self._process_word,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel
        }
    
    def process_file(self, file_path: str) -> List[Document]:
        """处理任意格式的文档"""
        file_path = Path(file_path)
        ext = file_path.suffix.lower()
        
        if ext not in self.supported_formats:
            raise ValueError(f"不支持的格式: {ext}")
        
        return self.supported_formats[ext](file_path)
    
    def _process_pdf(self, file_path: Path) -> List[Document]:
        """处理PDF文档"""
        from src.core.pdf_processor import PDFProcessor
        processor = PDFProcessor()
        return processor.process_pdf(str(file_path))
    
    def _process_text(self, file_path: Path) -> List[Document]:
        """处理文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return [Document(page_content=content, metadata={"source": str(file_path), "type": "text"})]
        except Exception as e:
            return [Document(page_content=f"文本处理失败: {str(e)}", metadata={"source": str(file_path), "type": "text"})]
    
    def _process_markdown(self, file_path: Path) -> List[Document]:
        """处理Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return [Document(page_content=content, metadata={"source": str(file_path), "type": "markdown"})]
        except Exception as e:
            return [Document(page_content=f"Markdown处理失败: {str(e)}", metadata={"source": str(file_path), "type": "markdown"})]
    
    def _process_word(self, file_path: Path) -> List[Document]:
        """处理Word/WPS文档"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            content = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))
            
            full_content = "\n".join(content)
            return [Document(page_content=full_content, metadata={"source": str(file_path), "type": "word"})]
            
        except ImportError:
            return [Document(page_content="需要安装python-docx库来处理Word文档", metadata={"source": str(file_path), "type": "word"})]
        except Exception as e:
            return [Document(page_content=f"Word处理失败: {str(e)}", metadata={"source": str(file_path), "type": "word"})]
    
    def _process_powerpoint(self, file_path: Path) -> List[Document]:
        """处理PowerPoint文档"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = [f"第{slide_idx + 1}页:"]
                
                # 提取标题
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"标题: {slide.shapes.title.text}")
                
                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                content.append("\n".join(slide_text))
            
            full_content = "\n\n".join(content)
            return [Document(page_content=full_content, metadata={"source": str(file_path), "type": "powerpoint"})]
            
        except ImportError:
            return [Document(page_content="需要安装python-pptx库来处理PowerPoint文档", metadata={"source": str(file_path), "type": "powerpoint"})]
        except Exception as e:
            return [Document(page_content=f"PowerPoint处理失败: {str(e)}", metadata={"source": str(file_path), "type": "powerpoint"})]
    
    def _process_excel(self, file_path: Path) -> List[Document]:
        """处理Excel文档"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_content = [f"工作表: {sheet_name}"]
                
                # 提取前100行数据
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if row_idx >= 100:  # 限制行数
                        break
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        sheet_content.append(" | ".join(row_data))
                
                content.append("\n".join(sheet_content))
            
            full_content = "\n\n".join(content)
            return [Document(page_content=full_content, metadata={"source": str(file_path), "type": "excel"})]
            
        except ImportError:
            return [Document(page_content="需要安装openpyxl库来处理Excel文档", metadata={"source": str(file_path), "type": "excel"})]
        except Exception as e:
            return [Document(page_content=f"Excel处理失败: {str(e)}", metadata={"source": str(file_path), "type": "excel"})]
    
    def process_image(self, file_path: Path) -> List[Document]:
        """处理图片（OCR）"""
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang='chi_sim+eng')
        
        return [Document(page_content=text, metadata={"source": str(file_path)})]
    
    def get_document_info(self, file_path: str) -> Dict[str, str]:
        """获取文档信息"""
        file_path = Path(file_path)
        stat = file_path.stat()
        
        return {
            'filename': file_path.name,
            'size': f"{stat.st_size / 1024:.1f} KB",
            'format': file_path.suffix.upper(),
            'pages': self._get_page_count(file_path)
        }
    
    def _get_page_count(self, file_path: Path) -> str:
        """获取页数或工作表数量"""
        ext = file_path.suffix.lower()
        
        if ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return str(len(reader.pages))
        elif ext == '.docx':
            doc = docx.Document(file_path)
            return str(len(doc.paragraphs) // 20 + 1)  # 估算
        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(file_path)
            return f"{len(wb.sheetnames)}个工作表"
        else:
            return "1"